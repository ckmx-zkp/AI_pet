# -*- coding: utf-8 -*-
"""AI Pet 正式版硬件电路设计说明 PPT 生成脚本"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

DARK = RGBColor(0x1F, 0x3B, 0x5C)     # 深蓝
ACCENT = RGBColor(0x2E, 0x86, 0xAB)   # 蓝绿
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)
GRAY = RGBColor(0x59, 0x59, 0x59)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xD6, 0x8A, 0x10)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set_ea(run):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
    rPr.append(ea)


def style(run, size=14, bold=False, color=GRAY):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    _set_ea(run)


def add_box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, level, space_after)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for (text, size, bold, color, level, space) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        if space:
            p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        style(r, size, bold, color)
    return tb


def new_slide(title=None, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    add_box(s, 0, 0, SW, Inches(0.16), fill=ACCENT)
    if title:
        add_text(s, Inches(0.55), Inches(0.32), Inches(12.3), Inches(0.9),
                 [(title, 26, True, DARK, 0, 0)])
        add_box(s, Inches(0.58), Inches(1.06), Inches(1.5), Pt(3), fill=ACCENT)
    if subtitle:
        add_text(s, Inches(2.25), Inches(0.95), Inches(10.6), Inches(0.4),
                 [(subtitle, 12, False, GRAY, 0, 0)])
    return s


def bullets(slide, x, y, w, h, items):
    """items: list of (text, level) 或 (text, level, bold, color)"""
    lines = []
    for it in items:
        text, level = it[0], it[1]
        bold = it[2] if len(it) > 2 else (level == 0)
        color = it[3] if len(it) > 3 else (DARK if level == 0 else GRAY)
        size = 15 if level == 0 else 13
        prefix = "▍ " if level == 0 else "· "
        lines.append((prefix + text if not text.startswith(("⚠", "🔴", "🟡", "✅", "❌")) else text,
                      size, bold, color, level, 6))
    add_text(slide, x, y, w, h, lines)


def add_table(slide, x, y, w, h, rows, col_widths=None, header_fill=DARK, font_size=12):
    from pptx.util import Cm
    shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                style(r, font_size, True, WHITE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                color = GRAY
                bold = False
                if isinstance(val, str) and val.startswith("🔴"):
                    color, bold = RED, True
                elif isinstance(val, str) and val.startswith("🟡"):
                    color, bold = ORANGE, True
                elif isinstance(val, str) and val.startswith(("✅", "🟢")):
                    color = GREEN
                style(r, font_size, bold, color)
    return tbl


# ============ S1 封面 ============
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=DARK)
add_box(s, 0, Inches(4.9), SW, Pt(2.5), fill=ACCENT)
add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(2.2), [
    ("AI Pet 正式版硬件", 44, True, WHITE, 0, 8),
    ("电路设计说明", 44, True, WHITE, 0, 0),
])
add_text(s, Inches(1), Inches(5.15), Inches(11.3), Inches(1.6), [
    ("ESP32-S3 主控  +  K230 视觉协处理器  +  ML307A 4G", 18, False, RGBColor(0xBF, 0xD7, 0xEA), 0, 10),
    ("依据：SCH_Schematic1_2026-08-02.pdf（嘉立创EDA，6页）/ 网表连通性核实 / xiaozhi-esp32 固件对照", 12, False, RGBColor(0x8F, 0xA8, 0xC0), 0, 4),
    ("2026-08-02", 12, False, RGBColor(0x8F, 0xA8, 0xC0), 0, 0),
])

# ============ S2 整机架构 ============
s = new_slide("整机架构", "三芯方案：P4 样机 → S3 正式版的预期切换")
arch = [
    ("ESP32-S3-WROOM-2\nN32R16V\n主控 · 跑小智固件", 0.7, 1.9, 3.2, 1.5, ACCENT),
    ("K230 核心模组\n140Pin 邮票孔\n视觉推理协处理", 5.05, 1.9, 3.2, 1.5, ACCENT),
    ("ML307A\n4G Cat.1\n数据通道", 9.4, 1.9, 3.2, 1.5, ACCENT),
    ("OV2640 DVP\n被动拍照", 0.7, 4.15, 1.5, 1.0, LIGHT),
    ("GC9A01 圆屏\n眼睛", 2.4, 4.15, 1.5, 1.0, LIGHT),
    ("ES7210+ES8311\n+NS4150B 音频", 4.1, 4.15, 1.9, 1.0, LIGHT),
    ("舵机 / 触摸\n/ 霍尔", 6.2, 4.15, 1.5, 1.0, LIGHT),
    ("CSI0 相机\n主动视觉", 7.9, 4.15, 1.5, 1.0, LIGHT),
    ("SIM + IPEX 天线", 9.6, 4.15, 1.5, 1.0, LIGHT),
    ("电池 + TP4054 充电", 11.3, 4.15, 1.5, 1.0, LIGHT),
]
for text, x, y, w, h, fill in arch:
    b = add_box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill,
                line=ACCENT if fill == LIGHT else None)
    tf = b.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        style(r, 12 if i == 0 else 10, i == 0, WHITE if fill == ACCENT else DARK)
add_text(s, Inches(0.7), Inches(3.55), Inches(12), Inches(0.5), [
    ("S3 ══ UART1 ══ K230        S3 ══ UART ══ ML307A        Type-C ══ K230 USB0（烧录/调试）", 13, True, GRAY, 0, 0)])
bullets(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.8), [
    ("方案定位：开发计划中「S3 + K230 单通路」降级备选的实际落地，主控降本约 30 元", 0),
    ("语音 / 表情 / 运动 / 联网由 S3 承担；K230 专职视觉推理，UART 上报结构化结果", 1),
    ("4G 仅作数据通道，ML307 自带音频全部悬空，语音链路统一在 S3 侧", 1),
])

# ============ S3 原理图页面地图 ============
s = new_slide("原理图概览", "6 页图纸 · 功能分区")
add_table(s, Inches(0.7), Inches(1.5), Inches(12), Inches(4.4), [
    ["页", "图纸标题", "功能内容", "关键器件"],
    ["P1", "主控，摄像，眼睛", "S3 主控、DVP 摄像头、GC9A01 眼睛屏、舵机接口、IO 扩展、CH340K 调试", "ESP32-S3-WROOM-2 / PCA9557 / CH340K / XC6206×2"],
    ["P2", "咪头，喇叭", "双模拟麦、四麦 ADC、音频 Codec、Class-D 功放、AEC 回采", "ES7210 / ES8311 / NS4150B / ZTS6216×2 / SY8088×2"],
    ["P3", "K230", "K230 核心模组、CSI0 摄像头接口、调试排针", "K230 模组 / XC6220-3.3V/1A"],
    ["P4", "4G模块", "ML307A、SIM 卡座、IPEX 天线、受控电源开关", "ML307A / SMN-305 / AO3415"],
    ["P5", "电池", "充电、4.0V 电源母线、4G 供电 DCDC", "TP4054 / TPS562200 / DSS24×2"],
    ["P6", "（预留）", "仅调试排针 P2", "—"],
], col_widths=[Inches(0.7), Inches(1.9), Inches(5.2), Inches(4.2)])

# ============ S4 语音链路 ============
s = new_slide("功能① 语音交互链路", "全图最成熟的部分，与小智参考设计同源")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.3), Inches(5.6), [
    ("拾音：2× ZTS6216 模拟麦，差分走线进 ES7210（0x41）", 0),
    ("MICBIAS12 供电，10µF+100nF 退耦，模拟前端规范", 1),
    ("AEC 硬件参考回采（亮点）", 0),
    ("NS4150B 输出 OUTP/OUTN 经 0Ω 回送 ES7210 第 3 通道", 1),
    ("对应固件 AUDIO_INPUT_REFERENCE=true，回声消除更稳", 1),
    ("放音：ES8311（0x18）→ NS4150B Class-D → J7 喇叭", 0),
    ("PA_EN 由 PCA9557 控制，10k 下拉保证上电静默", 1),
    ("I2S：MCLK=38 / BCK=14 / WS=13 / DI=12 / DO=45", 0),
    ("音频独立电源域", 0),
    ("SY8088×2 独立 buck 出 AU_3V3 / ADC_3V3（680k/150k→3.32V）", 1),
    ("AU_GND 经 0Ω 与主 GND 单点连接，数模分域正确", 1),
])
add_box(s, Inches(7.3), Inches(1.6), Inches(5.3), Inches(5.0), fill=LIGHT, line=ACCENT)
add_text(s, Inches(7.55), Inches(1.8), Inches(4.9), Inches(4.7), [
    ("信号流", 15, True, DARK, 0, 10),
    ("MIC ×2", 13, False, GRAY, 0, 2),
    ("   ↓ 差分", 11, False, GRAY, 0, 2),
    ("ES7210（4ch ADC，TDM）", 13, False, GRAY, 0, 2),
    ("   ↓ I2S_DI（51Ω 源端阻尼）", 11, False, GRAY, 0, 2),
    ("ESP32-S3（小智语音协议栈）", 13, True, DARK, 0, 2),
    ("   ↓ I2S_DO", 11, False, GRAY, 0, 2),
    ("ES8311（DAC/Codec）", 13, False, GRAY, 0, 2),
    ("   ↓ 差分 + 交流耦合", 11, False, GRAY, 0, 2),
    ("NS4150B（Class-D，5V）", 13, False, GRAY, 0, 2),
    ("   ↓", 11, False, GRAY, 0, 2),
    ("喇叭 ＋ 回采→ES7210 ch3", 13, True, ACCENT, 0, 8),
    ("结论：✅ 无需改动，直接复用固件双 codec 路径", 13, True, GREEN, 0, 0),
])

# ============ S5 视觉链路 ============
s = new_slide("功能② 视觉双通路", "主动视觉（K230）+ 被动拍照（S3），与开发计划一致")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.1), Inches(5.6), [
    ("通路 A：K230 主动视觉", 0),
    ("FPC1 22Pin → CSI0（2-lane MIPI），I2C0 带 2K 上拉", 1),
    ("XC6220-3.3V/1A 独立供电（3V3-EXT）", 1),
    ("推理结果经 UART1（S3 IO10/11）上报 face_pose / face_emotion", 1),
    ("Type-C 的 D+/D- 接 K230 USB0；UART0/3 引到 P2 排针调试", 1),
    ("通路 B：S3 被动拍照", 0),
    ("FPC3 24Pin DVP → OV2640（SCCB 0x60，共用 IO1/2 总线）", 1),
    ("XC6206-2.8V + XC6206-1.2V 双 LDO 供电，PWDN 走 PCA9557", 1),
    ("S3 无 ISP/JPEG 硬加速 → 只适合低频拍照上云，符合定位", 1),
])
add_box(s, Inches(7.1), Inches(1.6), Inches(5.5), Inches(5.0), fill=LIGHT, line=ACCENT)
add_text(s, Inches(7.35), Inches(1.8), Inches(5.1), Inches(4.7), [
    ("设计评价", 15, True, DARK, 0, 10),
    ("✅ 用比 P4 便宜的方案保住了两条视觉链路，", 13, False, GRAY, 0, 4),
    ("    是本版相对 P4 样机最成功的取舍", 13, False, GRAY, 0, 10),
    ("K230 模组上 DSI / TF / WiFi-SDIO / 触摸", 13, False, GRAY, 0, 4),
    ("全部悬空未用 —— K230 只做「相机进、结果出」", 13, False, GRAY, 0, 10),
    ("⚠ 主 3V3 由 K230 模组 VOUT_3V3 反供全板", 13, True, ORANGE, 0, 4),
    ("    K230 不可省配，且需核对其 3V3 输出电流余量", 13, False, GRAY, 0, 10),
    ("XCLK=IO5 PCLK=IO7 VSYNC=IO3 HREF=IO46", 12, False, ACCENT, 0, 2),
    ("D0~D7=IO16/18/17/8/15/6/4/9", 12, False, ACCENT, 0, 0),
])

# ============ S6 显示 ============
s = new_slide("功能③ 眼睛显示", "1.28 寸 GC9A01 圆屏 × 1")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.3), Inches(5.6), [
    ("J4 10Pin FPC：SPI 三线 SCK=41 / MOSI=40 / DC=39（无 MISO，正确）", 0),
    ("背光：3V3 → 10Ω → SI2301 P-MOS → BL_A", 0),
    ("IO42 经 1k 串阻控栅，10k 上拉保证上电灭屏 ✓", 1),
    ("待改问题", 0, True, ORANGE),
    ("LCD_CS 挂在 PCA9557（I2C）上：每次翻转都是一次 I2C 事务，", 1),
    ("拖慢刷屏且污染音频/摄像头共用的 I2C 总线", 1),
    ("LCD_RST 与系统 RESET 硬绑，固件无法独立复位屏幕", 1),
    ("改版方案（已评估 S3 引脚预算，见引脚页）", 0, True, ACCENT),
    ("LCD_CS → 直接接地（总线独占，模块本身无 CS 脚），固件 cs=-1", 1),
    ("LCD_RST → 挪到空闲的 IO19；IO20 留作备用", 1),
    ("空出的 PCA9557 IO0 顺手接 4G_PWR，解决 4G 电源控制悬空", 1),
])
add_box(s, Inches(7.3), Inches(1.6), Inches(5.3), Inches(5.0), fill=LIGHT, line=ACCENT)
add_text(s, Inches(7.55), Inches(1.8), Inches(4.9), Inches(4.7), [
    ("新板型显示配置（固件）", 15, True, DARK, 0, 10),
    ("SCK  = GPIO41", 13, False, GRAY, 0, 3),
    ("MOSI = GPIO40", 13, False, GRAY, 0, 3),
    ("DC   = GPIO39", 13, False, GRAY, 0, 3),
    ("RST  = GPIO19（改后）", 13, True, ACCENT, 0, 3),
    ("CS   = -1 / 接地（改后）", 13, True, ACCENT, 0, 3),
    ("BL   = GPIO42", 13, False, GRAY, 0, 12),
    ("功能缺口：全图仅 1 个屏接口", 13, True, ORANGE, 0, 4),
    ("进度看板规划的「第二只眼」无落点，", 13, False, GRAY, 0, 2),
    ("双眼版本需要改版预留", 13, False, GRAY, 0, 0),
])

# ============ S7 运动与感知 ============
s = new_slide("功能④ 运动与感知", "舵机 + 触摸/霍尔")
bullets(s, Inches(0.7), Inches(1.45), Inches(12), Inches(5.6), [
    ("舵机（CN5，MX3.0 2x2）", 0),
    ("信号 SE_PWMO = IO23 直驱（3.3V 电平对 5V 舵机一般可识别）", 1),
    ("电源 5V 经 D5（DSS24，额定 2A）供给，图纸标注堵转电流 2.2A", 1),
    ("🔴 D5 在堵转时已超限：建议换 ≥3A 肖特基或直短", 2, True, RED),
    ("图纸批注「到位后关闭 PWM 省电」方向正确，固件实现即可", 1),
    ("触摸 / 霍尔（P1，PH2.0-6P）", 0),
    ("HE_EN + CM_1~4 共 5 路，全部经 PCA9557 采集", 1),
    ("低速开关型信号放 I2C 扩展器是合理用法 ✓", 1),
    ("PCA9557 地址 0x19，与 ES8311(0x18) / ES7210(0x41) / OV2640(0x60) 无冲突 ✓", 1),
    ("功能缺口（对比固件 V0.3 规划）", 0, True, ORANGE),
    ("WS2812 灯带：无电路、无连接器 ❌", 1),
    ("双舵机：仅 1 路接口 ❌（第二只眼同样无落点）", 1),
])

# ============ S8 通信 ============
s = new_slide("功能⑤ 通信", "WiFi/BT 内置 + ML307A 4G")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.2), Inches(5.6), [
    ("WiFi / BT：S3-WROOM-2 模组内置，免射频设计（降本主力）", 0),
    ("4G：ML307A + SMN-305 SIM 座 + IPEX 天线", 0),
    ("SIM 带 PESD 静电防护、SIM_DET 下拉 ✓；π 型匹配 NC 预留 ✓", 1),
    ("与 S3 走 UART（IO47/48），固件复用上游 Ml307Board", 1),
    ("PWR_ON/OFF 4.7k 下拉，RESET/BOOT_MODE 引到测试排针 ✓", 1),
    ("ML307 自带音频（MIC/SPK）全部悬空：4G 只作数据通道 ✓", 1),
    ("VBAT：3×100µF + 多档电容，符合 ML307 硬件手册 ✓", 1),
    ("已知问题", 0, True, ORANGE),
    ("4G_PWR 悬空：AO3415 栅极被上拉默认关断，4G 无法上电", 1, False, RED),
    ("建议挂到改版后空出的 PCA9557 IO0（慢速开关正合适）", 1),
])
add_box(s, Inches(7.1), Inches(1.6), Inches(5.5), Inches(5.0), fill=LIGHT, line=ACCENT)
add_text(s, Inches(7.35), Inches(1.8), Inches(5.1), Inches(4.7), [
    ("USB-C 与烧录通道", 15, True, DARK, 0, 10),
    ("Type-C 的 D+/D- → K230 USB0", 13, False, GRAY, 0, 4),
    ("S3 烧录：CH340K + DTR/RTS 自动下载", 13, False, GRAY, 0, 4),
    ("电路（小智全系标准方案），走 H1 排针", 13, False, GRAY, 0, 4),
    ("S3 原生 USB（IO19/20）悬空 → 已改作 LCD 用", 13, False, GRAY, 0, 12),
    ("⚠ CC 电阻 R8/R15 = 10kΩ", 13, True, ORANGE, 0, 4),
    ("UFP 规范应为 5.1kΩ Rd；部分 PD 充电器", 13, False, GRAY, 0, 2),
    ("会因此不放电 —— 舵机/4G 吃电流，建议必改", 13, False, GRAY, 0, 12),
    ("调试：K230 UART0/3 → P2 排针；S3 UART0 → CH340K", 12, False, ACCENT, 0, 0),
])

# ============ S9 电源系统 ============
s = new_slide("功能⑥ 电源系统", "全图最大短板：电池放电路径断裂")
add_box(s, Inches(0.7), Inches(1.5), Inches(7.2), Inches(4.3), fill=LIGHT, line=ACCENT)
add_text(s, Inches(0.95), Inches(1.65), Inches(6.8), Inches(4.1), [
    ("电源树", 15, True, DARK, 0, 8),
    ("USB-C 5V ─┬─ K230 VIN → 模组 VOUT_3V3 → 全板 3V3（S3/屏/外设）", 12, False, GRAY, 0, 3),
    ("              ├─ XC6220 → 3V3-EXT（K230 相机）", 12, False, GRAY, 0, 3),
    ("              ├─ D5 → 舵机 5V", 12, False, GRAY, 0, 3),
    ("              ├─ NS4150B 功放 5V", 12, False, GRAY, 0, 3),
    ("              ├─ TP4054 → 电池 4.2V/2000mAh（300mA ✓）", 12, False, GRAY, 0, 3),
    ("              └─ TPS562200 → +4.3V ─D4→ V4.0V ─┬─ AO3415 → ML307", 12, False, GRAY, 0, 3),
    ("电池 ─D3────────────────→ V4.0V ─┘            └─ SY8088×2 → 音频 3V3", 12, False, GRAY, 0, 8),
    ("5V 只来自 USB；主 3V3 由 K230 模组反供", 13, True, RED, 0, 3),
    ("→ 拔掉 USB：K230 停 → 3V3 停 → S3/屏/舵机全灭", 13, True, RED, 0, 3),
    ("→ 电池只能维持 4G 与音频 buck，主控已死，功能不成立", 13, True, RED, 0, 0),
])
bullets(s, Inches(8.2), Inches(1.5), Inches(4.5), Inches(4.4), [
    ("修复方向", 0),
    ("补 battery→5V boost，或 V4.0V 直供 K230 + 独立 3V3 稳压器", 1),
    ("投板前必须修", 1, True, RED),
    ("5V 预算", 0),
    ("舵机堵转 2.2A + K230 峰值 ~1.5A + 充电 0.3A + 功放 ~0.6A", 1),
    ("适配器需 5V/4A 级；CC 阻值错误会让 PD 头限流", 1),
    ("其他", 0),
    ("NTC 测温电路正确，但采样未接到任何 ADC（单点网络）", 1),
    ("TP4054 无均充路径：系统走 5V、电池只充电，等效回避 ✓", 1),
])

# ============ S10 S3 引脚预算 ============
s = new_slide("ESP32-S3 引脚预算", "WROOM-2-N32R16V · 41Pin 全部核对")
add_table(s, Inches(0.7), Inches(1.5), Inches(12), Inches(4.6), [
    ["功能", "GPIO", "说明"],
    ["I2C 总线", "IO1(SDA) / IO2(SCL)", "ES8311 0x18 / PCA9557 0x19 / ES7210 0x41 / OV2640 0x60"],
    ["DVP 摄像头", "IO3~IO9 / IO15~IO18 / IO46", "XCLK=5 PCLK=7 VSYNC=3 HREF=46 + 8 数据线"],
    ["I2S 音频", "IO12 / IO13 / IO14 / IO38 / IO45", "DI / WS / BCK / MCLK / DO"],
    ["UART1 → K230", "IO10 / IO11", "并引出到舵机口 CN5"],
    ["UART → ML307", "IO47 / IO48", "4G 数据通道"],
    ["UART0 调试", "IO43 / IO44", "CH340K 自动下载"],
    ["LCD", "IO39 / IO40 / IO41 / IO42", "DC / MOSI / SCK / 背光"],
    ["舵机 PWM", "IO23", "SE_PWMO"],
    ["BOOT", "IO0", "按键 + 自动下载"],
    ["PSRAM 占用", "IO35 / IO36 / IO37", "八线 PSRAM，不可用（网表悬空即此因）"],
    ["✅ 空闲可用", "IO19 / IO20", "原生 USB 脚，改作 LCD_RST / 备用"],
], col_widths=[Inches(2.6), Inches(3.4), Inches(6.0)], font_size=12)

# ============ S11 核心器件清单 ============
s = new_slide("核心器件清单", "正式版 BOM 关键件")
add_table(s, Inches(0.7), Inches(1.45), Inches(12), Inches(5.2), [
    ["位号", "器件", "作用域", "备注"],
    ["U3", "ESP32-S3-WROOM-2-N32R16V", "主控", "32MB Flash + 16MB 八线 PSRAM"],
    ["M1", "K230 核心模组（140Pin）", "视觉协处理", "板载 DCDC 反供全板 3V3，不可省配"],
    ["U28", "ML307A 4G Cat.1（LGA-94）", "通信", "仅数据通道；配 SMN-305 SIM 座"],
    ["U10", "ES7210", "音频 ADC", "四通道，ch3 作 AEC 参考回采"],
    ["U9", "ES8311", "音频 Codec", "I2C 0x18"],
    ["U11", "NS4150B", "Class-D 功放", "5V 供电，PA_EN 受控"],
    ["U2", "PCA9557PW", "IO 扩展", "0x19：LCD_CS/PA_EN/DVP_PWDN/触摸/霍尔"],
    ["U8", "CH340K", "调试", "S3 烧录，DTR/RTS 自动下载"],
    ["U7", "TP4054", "充电", "300mA → 4.2V/2000mAh"],
    ["U14", "TPS562200", "DCDC", "5V→4.3V，供 4G 与音频 buck"],
    ["U5/U6", "SY8088AAC ×2", "音频 LDO 域", "AU_3V3 / ADC_3V3 独立供电"],
    ["U1/U4/U12", "XC6220 / XC6206 ×2", "相机电源", "3V3-EXT / CAM_2V8 / CAM_1V2"],
    ["—", "GC9A01 1.28 寸圆屏", "眼睛", "SPI，资料见 模块-7针蓝板模块"],
    ["—", "OV2640 / ZTS6216×2", "相机 / 麦克风", "DVP 拍照 / 差分模拟麦"],
], col_widths=[Inches(1.3), Inches(3.3), Inches(2.0), Inches(5.4)], font_size=11)

# ============ S12 问题清单 ============
s = new_slide("问题清单", "按严重程度分级（网表连通性核实）")
add_table(s, Inches(0.7), Inches(1.45), Inches(12), Inches(5.2), [
    ["级别", "问题", "影响", "建议"],
    ["🔴", "电池放电路径断裂：5V 仅来自 USB，主 3V3 由 K230 反供", "拔掉 USB 整机断电，移动使用功能不成立", "补 boost 或 V4.0V 直供 K230+独立 3V3，投板前必改"],
    ["🔴", "4G_PWR 控制网络悬空（仅连 R65）", "AO3415 默认关断，ML307 无法上电", "挂到改版空出的 PCA9557 IO0"],
    ["🟡", "舵机 5V 串联 D5（DSS24，2A），堵转 2.2A", "堵转超二极管额定值", "换 ≥3A 肖特基或直短"],
    ["🟡", "LCD_CS 挂 I2C 扩展器；LCD_RST 绑系统 RESET", "刷屏慢、污染 I2C 总线、无法独立复位", "CS 接地、RST 挪 IO19"],
    ["🟡", "BAT__10K_ADC / VCC3V3 单点网络", "电池 NTC 测温功能断路", "接到 S3 ADC 引脚与 3V3"],
    ["🟡", "USB CC 电阻 10kΩ（应为 5.1kΩ Rd）", "部分 PD 充电器不放电/限流", "改 5.1kΩ"],
    ["🟢", "S3 烧录走 H1 排针；原生 USB 悬空", "用户烧录不便（产线可接受）", "维持现状，IO19/20 已改作 LCD"],
    ["🟢", "网表元件值与原理图不一致（如 R58）", "BOM 可能出错", "以原理图 PDF 为准重新导出网表"],
    ["🟢", "Schematic111 网表损坏（`.` 被替换为 123）", "文件不可用", "删除或重新导出"],
], col_widths=[Inches(0.8), Inches(4.3), Inches(3.3), Inches(3.6)], font_size=11)

# ============ S13 改版建议 ============
s = new_slide("改版建议", "一次改版解决问题 + 为 V0.3 预留")
bullets(s, Inches(0.7), Inches(1.45), Inches(12), Inches(4.0), [
    ("必修（投板前）", 0, True, RED),
    ("电源：打通电池→系统放电路径（boost 或 K230 改由 V4.0V 供电 + 独立 3V3）", 1),
    ("4G_PWR 接 PCA9557 IO0；CC 电阻改 5.1kΩ；D5 换 ≥3A", 1),
    ("LCD_CS 接地、LCD_RST 挪 IO19；NTC 采样接 S3 ADC", 1),
    ("建议同期预留（避免 V0.3 再次改板）", 0, True, ORANGE),
    ("第二只眼接口（共用 SPI，独立 CS/BL）", 1),
    ("WS2812 灯带接口（1 个 GPIO + 5V，注意 5V 预算）", 1),
    ("第二路舵机接口（信号 + 独立供电考量）", 1),
    ("固件配套", 0, True, ACCENT),
    ("新建 board（复用 esp32-p4-ai-pet）：ES8311+ES7210 双 codec、GC9A01、", 1),
    ("OV2640 DVP、Ml307Board、PCA9557 驱动、K230 UART 链路（main/pet 类型层已有设计稿）", 1),
])

# ============ S14 总结 ============
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=DARK)
add_box(s, 0, Inches(1.9), SW, Pt(2.5), fill=ACCENT)
add_text(s, Inches(1), Inches(0.8), Inches(11.3), Inches(1.0), [
    ("总结", 36, True, WHITE, 0, 0)])
add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(4.4), [
    ("语音、视觉、4G 三条主功能链路与小智固件资产高度对齐，是合格的量产雏形。", 18, False, WHITE, 0, 14),
    ("P4 → S3 切换符合预期：降本约 30 元，同时保住了 K230 主动视觉 + S3 被动拍照双链路。", 18, False, WHITE, 0, 14),
    ("两处硬伤（电池放电路径、4G_PWR 悬空）必须在投板前修复；", 18, True, RGBColor(0xF5, 0xB7, 0xB1), 0, 6),
    ("五处中等问题建议同版一并处理，并为双眼 / 灯带 / 双舵机预留接口。", 18, True, RGBColor(0xF5, 0xB7, 0xB1), 0, 0),
])

prs.save("AI_Pet正式版硬件电路设计说明.pptx")
print("saved:", "AI_Pet正式版硬件电路设计说明.pptx")
