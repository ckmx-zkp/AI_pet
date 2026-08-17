# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = r"D:\Home_Work\hardware\AI_Pet_ESP32-S3_Hardware_Audit_Camera_P0_2026-08-08.pptx"
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]

INK = RGBColor(20, 31, 46)
SLATE = RGBColor(73, 86, 104)
PANEL = RGBColor(239, 243, 247)
RULE = RGBColor(194, 203, 213)
BLUE = RGBColor(43, 119, 196)
CYAN = RGBColor(86, 184, 220)
RED = RGBColor(194, 54, 51)
AMBER = RGBColor(195, 124, 24)
GREEN = RGBColor(28, 135, 90)
WHITE = RGBColor(255, 255, 255)
FONT = "Microsoft YaHei"

def set_font(run, size, color=INK, bold=False):
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr(); rpr.set(qn('a:ea'), FONT)

def rect(slide, x, y, w, h, fill=WHITE, line=RULE, radius=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(0.8)
    return s

def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=t.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; set_font(r,size,color,bold)
    return t

def bullets(slide, x, y, w, items, size=17, color=SLATE, h=5.1):
    t=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=t.text_frame; tf.clear(); tf.word_wrap=True
    for i,(txt,clr,bold) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text="• "+txt; p.space_after=Pt(10)
        for r in p.runs: set_font(r,size,clr or color,bold)
    return t

def base(title, subtitle, page):
    s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
    rect(s,0,0,13.333,.12,BLUE,BLUE)
    textbox(s,.65,.36,11.8,.52,title,34,INK,True)
    textbox(s,.67,.95,11.2,.28,subtitle,14,SLATE)
    textbox(s,.67,7.08,10.8,.2,"来源：SCH_Schematic1_2026-08-08.pdf、Netlist_Schematic1_2026-08-08.tel、xiaozhi-esp32 代码（只读审核）",9,SLATE)
    textbox(s,12.15,7.03,.5,.25,f"{page:02d}",10,SLATE,True,PP_ALIGN.RIGHT)
    return s

def label(slide,x,y,w,text,fill=INK):
    r=rect(slide,x,y,w,.32,fill,fill,True); textbox(slide,x+.08,y+.045,w-.16,.18,text,10,WHITE,True,PP_ALIGN.CENTER)
    return r

# 1 cover
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=INK
rect(s,0,4.95,13.333,2.55,BLUE,BLUE)
textbox(s,.85,1.5,11.8,.7,"AI Pet ESP32‑S3 硬件原理图",42,WHITE,True)
textbox(s,.85,2.3,11.8,.65,"嵌入式软件工程师审核",42,WHITE,True)
textbox(s,.88,5.35,8.8,.35,"投板就绪度 · 烧录与驱动闭环 · 双 SPI 眼屏",20,WHITE,False)
textbox(s,.88,6.25,8.8,.28,"审核范围：S3 主控、下载、音频、双眼、OV2640、I2C/I2S、4G 与软件 bring-up",13,RGBColor(211,230,244))
textbox(s,.88,6.72,8.8,.22,"2026-08-08 原理图版本 · 只读审阅",11,RGBColor(211,230,244))

# 2 verdict
s=base("结论：可做 EVT，但当前不满足正式投板 Go", "S3 选型与双眼拓扑通过；三个 P0 闭环问题必须先修订",2)
rect(s,.72,1.55,4.0,4.75,RGBColor(253,241,240),RED)
textbox(s,1.0,1.88,3.4,.55,"NO‑GO（当前）",26,RED,True,PP_ALIGN.CENTER)
textbox(s,1.0,2.65,3.35,2.8,"不建议按当前原理图作为正式功能板下单。\n\nS3 核心能力可实现；但下载口、功放供电、4G UART 电平尚未形成稳定的软硬件闭环。",20,INK,False,PP_ALIGN.CENTER)
rect(s,4.95,1.55,7.65,4.75,PANEL,RULE)
label(s,5.25,1.85,1.6,"可保留",GREEN)
bullets(s,5.25,2.28,6.85,[
 ("ESP32‑S3 + Wi‑Fi/BT + OTA + 小智语音",GREEN,True),
 ("ES7210 / ES8311 音频架构与 I2S 引脚资源",GREEN,True),
 ("GC9A01 双屏：共享 SPI、双 CS 的硬件拓扑",GREEN,True),
 ("OV2640：首板 P0 功能，采用 S3 DVP 单帧 JPEG",GREEN,True),
 ("K230 作为 UART 视觉事件源的后续扩展",GREEN,True),
],18)
label(s,5.25,4.78,1.6,"投板前修复",RED)
bullets(s,5.25,5.2,6.85,[
 ("Type‑C / CH340 / S3 下载路径必须明确",RED,True),
 ("NS4150B 的 VBUS 与系统 5V 供电关系必须闭环",RED,True),
 ("若装 4G，必须确认 1.8V UART 并补电平转换",RED,True),
],17,h=1.35)

# 3 capability
s=base("S3 的产品能力覆盖充分，P4→S3 是预期收敛", "产品是人格陪伴电子宠；核心是语音、双眼与表达，不是 MIPI 视觉平台",3)
rows=[("能力","S3 判断","审核理由"),("语音 / OTA / Wi‑Fi","通过","S3 原生 Wi‑Fi/BT，现有工程已支持 ESP32‑S3 target"),("双麦 + 扬声器","通过","I2S + 双 codec 的 IO 分配无直接冲突"),("双眼 / MCP 表情","通过","GC9A01 的 SPI 多从设备模式适配"),("OV2640 首板拍照","通过","S3 DVP + esp32_camera；先验收 QVGA JPEG，再上 VGA"),("K230 主动视觉","可行","S3 只收 UART 结构化结果，不承担推理"),("MIPI‑CSI / DSI","不保留","属于 P4 路线，非当前产品核心")]
table=s.shapes.add_table(len(rows),3,Inches(.72),Inches(1.55),Inches(11.9),Inches(4.95)).table
table.columns[0].width=Inches(2.55); table.columns[1].width=Inches(1.55); table.columns[2].width=Inches(7.8)
for r,row in enumerate(rows):
    for c,v in enumerate(row):
        cell=table.cell(r,c); cell.fill.solid(); cell.fill.fore_color.rgb=INK if r==0 else (WHITE if r%2 else PANEL)
        cell.margin_left=cell.margin_right=Inches(.1); cell.margin_top=cell.margin_bottom=Inches(.06); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.text=str(v); p=cell.text_frame.paragraphs[0]
        for rr in p.runs: set_font(rr,14,WHITE if r==0 else (GREEN if v=="通过" else (AMBER if v=="可行" else SLATE)),r==0 or c==1)
textbox(s,.78,6.56,11.7,.28,"软件迁移重点：新建 esp32‑s3‑ai‑pet 板型；迁移眼睛状态机与 MCP；替换 P4 CSI 与单屏输出层。",16,INK,True)

# 4 pinmap
s=base("S3 引脚分配可驱动全部首版外设", "关键功能组之间未发现 GPIO 复用冲突；GPIO19/20 的取舍需单独确认",4)
cols=[(.75,"总线 / IO","GPIO","设备"), (4.95,"显示","39 / 40 / 41","DC / MOSI / SCK"), (9.0,"音频","12 / 13 / 14 / 38 / 45","DI / WS / BCK / MCLK / DO")]
for x,h,a,b in cols:
    textbox(s,x,1.5,3.3,.35,h,20,INK,True)
    rect(s,x,1.95,3.3,1.0,PANEL,RULE); textbox(s,x+.18,2.14,2.9,.26,a,19,BLUE,True); textbox(s,x+.18,2.54,2.9,.22,b,14,SLATE)
for x,h,a,b in [(.75,"I2C","1 / 2","ES8311、ES7210、PCA9539、OV2640"),(4.95,"双眼片选","19 / 20","左眼 CS / 右眼 CS"),(9.0,"摄像头 DVP","3–9、15–18、46","OV2640 数据与同步")]:
    textbox(s,x,3.25,3.3,.35,h,20,INK,True); rect(s,x,3.7,3.3,1.0,PANEL,RULE); textbox(s,x+.18,3.89,2.9,.26,a,19,BLUE,True); textbox(s,x+.18,4.29,2.9,.22,b,14,SLATE)
for x,h,a,b in [(.75,"外设","23","舵机 PWM"),(4.95,"4G UART","47 / 48","ML307A"),(9.0,"下载 UART","U0TXD / U0RXD","CH340 自动下载")]:
    textbox(s,x,5.0,3.3,.35,h,20,INK,True); rect(s,x,5.45,3.3,.75,PANEL,RULE); textbox(s,x+.18,5.6,2.9,.26,a,19,BLUE,True); textbox(s,x+.18,5.94,2.9,.22,b,14,SLATE)
textbox(s,.78,6.55,11.6,.3,"注意：GPIO19/20 同时是 S3 原生 USB D−/D+。眼屏占用它们本身可行，但下载必须稳定走 CH340，不能再依赖 S3 原生 USB。",15,AMBER,True)

# 5 spi
s=base("新双 SPI 眼屏：硬件拓扑正确，软件输出层需重写", "两块 GC9A01 是 SPI 多从机；一条总线、两个 panel、单一动画状态机",5)
# connectors first
for y in [2.22,2.85,3.48,4.11]:
    line=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(3.3),Inches(y),Inches(9.9),Inches(y)); line.line.color.rgb=RULE; line.line.width=Pt(1.4)
rect(s,.75,1.72,2.55,3.0,INK,INK); textbox(s,1.0,2.1,2.05,.42,"ESP32‑S3",25,WHITE,True,PP_ALIGN.CENTER); textbox(s,1.0,2.75,2.05,1.2,"GPIO41 SCK\nGPIO40 MOSI\nGPIO39 DC\nRESET",16,WHITE,False,PP_ALIGN.CENTER)
rect(s,9.9,1.63,2.35,1.42,PANEL,BLUE); textbox(s,10.15,1.92,1.85,.25,"左眼 GC9A01",18,INK,True,PP_ALIGN.CENTER); textbox(s,10.15,2.32,1.85,.22,"CS = GPIO19",14,BLUE,True,PP_ALIGN.CENTER)
rect(s,9.9,3.8,2.35,1.42,PANEL,BLUE); textbox(s,10.15,4.09,1.85,.25,"右眼 GC9A01",18,INK,True,PP_ALIGN.CENTER); textbox(s,10.15,4.49,1.85,.22,"CS = GPIO20",14,BLUE,True,PP_ALIGN.CENTER)
textbox(s,3.62,1.95,5.7,.4,"SCK / MOSI / DC / RESET 共用",20,INK,True,PP_ALIGN.CENTER)
textbox(s,3.7,5.25,5.6,.42,"每帧：渲染 → RGB565 字节交换一次 → 左 panel → 右 panel",18,BLUE,True,PP_ALIGN.CENTER)
textbox(s,.85,6.13,11.6,.36,"禁止将现有 FlushFrame() 对左右屏各执行一次：它会原地交换 framebuffer，第二屏颜色会被换回。应仅有一个动画任务与一份同步状态。",15,RED,True)

# 6 P0
s=base("P0：三个问题会直接造成烧录或功能失效", "这些问题不能用软件重试规避，需在投板前由硬件闭环",6)
issues=[("01","Type‑C 下载路径不闭环","Type‑C D+/D−仍在 P4 残留 USB 网络；CH340 走 H1。\n需定义唯一的“Type‑C → CH340 → S3 UART0”路径。"),("02","功放供电网需确认","NS4150B VCC 接 VBUS；系统主电源另名 5V。\n若两网不相连，电池/DC 供电时 TTS 将无声。"),("03","4G UART 电平未关闭","图中已有“确认都是 1.8V电平？”注记。\n装配 ML307A 前必须确认并加双向电平转换。")]
for i,(num,title,body) in enumerate(issues):
    x=.75+i*4.12; rect(s,x,1.62,3.75,4.72,RGBColor(253,241,240),RED); label(s,x+.24,1.92,.56,num,RED); textbox(s,x+.24,2.53,3.2,.5,title,21,INK,True); textbox(s,x+.24,3.28,3.12,1.9,body,16,SLATE); textbox(s,x+.24,5.7,3.0,.22,"投板前必须关闭",14,RED,True)

# 7 P1
s=base("P1：这些风险会导致“能启动但不稳定”", "建议与 P0 同版修订，避免把硬件不确定性带入固件调试",7)
items=[("I2C 上拉与地址","1k 与 4.7k 疑似并联；建议统一 2.2k–4.7k。\n地址统一标成 7-bit，避免驱动扫描失败。",AMBER),("I2S 负载","MCLK / BCLK / WS 各约 44pF 到地。\n建议预留 DNP 或降低容值，避免时钟边沿裕量不足。",AMBER),("屏幕复位","双眼 RESET 与系统 EN 共用。\n首板可接受，但软件无法对单块故障屏做硬复位。",AMBER),("测试可达性","现有测试点集中在 4G USB。\n补 5V、3V3、4V、EN、IO0、UART、I2C、SPI 测试点。",AMBER)]
for i,(t,b,c) in enumerate(items):
    x=.75+(i%2)*6.1; y=1.62+(i//2)*2.35; rect(s,x,y,5.7,1.78,PANEL,RULE); label(s,x+.24,y+.2,1.1,"P1",AMBER); textbox(s,x+1.55,y+.22,3.85,.3,t,20,INK,True); textbox(s,x+.25,y+.72,5.05,.7,b,16,SLATE)

# 8 bringup
s=base("OV2640：首板 P0 功能，S3 DVP 引脚无冲突", "实现关键是地址、上电时序、FPC 对应关系与 DVP 走线，不是保留 P4 CSI",8)
for x,w,title,body,col in [
 (.75,3.75,"硬件映射","SCCB：GPIO1 / GPIO2\nXCLK=GPIO5；PCLK=GPIO7\nVSYNC=GPIO3；HREF=GPIO46\nD0–D7：GPIO16/18/8/17/15/6/4/9",BLUE),
 (4.8,3.75,"软件顺序","1  初始化 I2C / PCA9539\n2  PWDN 保持高（休眠）\n3  PWDN 拉低、启动 XCLK\n4  SCCB 探测 7-bit 地址 0x30\n5  QVGA JPEG 单帧采集，再上 VGA",GREEN),
 (8.85,3.75,"投板前闭环","确认 OV2640 FPC 管脚序与 D0–D7 对应\n核对 0x60 是否为 8-bit 写地址\nDVP 短线、完整地参考；预留 0/22R\n增加 CAM_2V8 / CAM_1V2 / SCCB / XCLK 测试点",RED)]:
    rect(s,x,1.62,w,4.72,PANEL,col); label(s,x+.22,1.9,1.25,"首板必测",col); textbox(s,x+.24,2.48,w-.45,.38,title,21,INK,True); textbox(s,x+.24,3.12,w-.5,2.6,body,15,SLATE)
textbox(s,.85,6.65,11.45,.28,"首板验收：可重复上电识别 OV2640、连续抓取 JPEG 无花屏/超时；仅按需拍照，不做连续 RGB565 视频流。",15,GREEN,True)

# 9 bringup
s=base("建议的投板门槛与软件 Bring‑up 验收顺序", "摄像头纳入 EVT 核心验收；4G/K230 不应阻塞语音、双眼与拍照闭环",9)
steps=[("1","下载","Type‑C/CH340、EN、IO0：稳定擦写、串口日志"),("2","电源","5V / 3V3 / 4V 上电、负载、复位时序"),("3","总线","I2C 扫描与 ES8311/ES7210/PCA9539 地址核验"),("4","摄像头","PWDN → XCLK → SCCB 0x30 → QVGA JPEG 单帧"),("5","音频","双麦录音、I2S 时钟、播放、AEC 回采"),("6","双眼","左/右单测 → 同步表情、眨眼、睡眠、背光"),("7","外设","舵机 PWM、触摸/霍尔；WS2812 需补电路后纳入"),("8","后续","ML307A、K230 UART 分阶段启用")]
for i,(n,t,b) in enumerate(steps):
    y=1.42+i*.62; rect(s,.85,y,.52,.43,BLUE,BLUE,True); textbox(s,.85,y+.075,.52,.17,n,14,WHITE,True,PP_ALIGN.CENTER); textbox(s,1.62,y+.01,1.1,.27,t,17,INK,True); textbox(s,2.78,y+.03,9.55,.27,b,15,SLATE)
textbox(s,.85,6.55,11.45,.32,"Go 判据：下载稳定 + OV2640 JPEG 拍照稳定 + 音频双向稳定 + 两屏同步刷新 + 静态功耗与异常复位可复现。",15,GREEN,True)

# 10 close
s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=INK
textbox(s,.85,1.25,11.6,.65,"建议：关闭三项 P0，并完成摄像头验证后投板",33,WHITE,True)
textbox(s,.88,2.25,10.8,.38,"S3 与新 SPI 屏不是风险来源；它们是当前产品定位最匹配的硬件基础。",22,RGBColor(210,229,244))
rect(s,.88,3.25,11.55,.9,BLUE,BLUE); textbox(s,1.15,3.52,11.0,.3,"投板目标应是：稳定对话 · 双眼有生命感 · 易烧录调试 · 可定位故障",23,WHITE,True,PP_ALIGN.CENTER)
textbox(s,.88,5.3,10.8,.35,"正式版 Go 条件",18,CYAN,True)
bullets(s,.9,5.72,10.8,[("下载口、功放供电、4G 电平三项 P0 关闭",WHITE,True),("OV2640 可稳定完成 QVGA JPEG 单帧拍照；双屏 S3 板型完成编译",WHITE,True),("PCB / Gerber / BOM / DRC / 测试点审查完成",WHITE,True)],17,WHITE,h=1.08)

prs.save(OUT)
print(OUT)
